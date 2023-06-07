#https://pypi.org/project/PyPDF2/

import os
import re
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document

#.docx .pptx .pdf
def obtener_nombre_archivos(extension):
    archivos = []
    for archivo in os.listdir("./dataset"):
        if archivo.endswith(extension):
            archivos.append(archivo)
    return archivos

def leer_pdf(ruta_archivo):
    archivoLectura = open(ruta_archivo,'rb')
    documento = PdfReader(archivoLectura)
    texto_documento = ""

    for pagina in documento.pages:
        texto_documento += pagina.extract_text()

    texto_limpio = re.sub(r'\W+', ' ', texto_documento)
    #texto_limpio = re.sub(r'\s+', ' ', texto_limpio)
    return {"ruta": ruta_archivo,"texto": texto_limpio}

def leer_pptx(filepath):
    presentation = Presentation(filepath)
    texto_documento = ""

    # Recorrer todas las diapositivas
    for slide in presentation.slides:
        # Recorrer todos los elementos de texto en la diapositiva
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        texto_documento += run.text
                        texto_documento += " "
                        
    texto_limpio = re.sub(r'\W+', ' ', texto_documento)
    texto_limpio = re.sub(r'\s+', ' ', texto_limpio)
    return texto_limpio

def leer_docx(ruta_archivo):
    doc = Document(ruta_archivo)
    texto_documento = ""

    # Recorrer todos los párrafos del documento
    for paragraph in doc.paragraphs:
        texto_documento += paragraph.text
        texto_documento += " "

    texto_limpio = re.sub(r'\W+', ' ', texto_documento)
    texto_limpio = re.sub(r'\s+', ' ', texto_limpio)
    return texto_limpio

#Ejecucion
nombres_archivos_pdf = obtener_nombre_archivos(".pdf")
nombres_archivos_docx = obtener_nombre_archivos(".docx")
nombres_archivos_pptx = obtener_nombre_archivos(".pptx")
page = leer_pdf("./dataset/MKT 2016 - Alan Szpigiel - TP4 (1).pdf")
power = leer_pptx("./dataset/Domótica_Final.pptx.pptx")
docx = leer_docx("./dataset/TP1.docx")
#print(page["texto"])


import spacy
nlp = spacy.load("es_core_news_sm")
nlp.pipe_names